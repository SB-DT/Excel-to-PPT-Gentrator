#!/usr/bin/env python
# coding: utf-8

# In[ ]:


import tkinter as tk
import customtkinter as ctk
from tkinter import filedialog, messagebox
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import webbrowser

# Set appearance mode and default color theme
ctk.set_appearance_mode("Dark")  # Set to Dark for a modern look
ctk.set_default_color_theme("dark-blue")  # Choose a custom theme

class App(ctk.CTk):
    def __init__(self):
        super().__init__()

        # Configure window
        self.title("Excel to PPT Generator")
        self.geometry("1200x600")

        # Configure grid layout (4x4)
        self.grid_columnconfigure(1, weight=1)
        self.grid_columnconfigure((2, 3), weight=0)
        self.grid_rowconfigure((0, 1, 2, 3), weight=1)

        # Create sidebar frame with widgets
        self.sidebar_frame = ctk.CTkFrame(self, width=250, corner_radius=10)
        self.sidebar_frame.grid(row=0, column=0, rowspan=4, sticky="ns")
        self.sidebar_frame.grid_rowconfigure(9, weight=1)
        
        # Title label
        self.logo_label = ctk.CTkLabel(
            self.sidebar_frame, 
            text="Excel to PPT Generator", 
            font=ctk.CTkFont(size=24, weight="bold")
        )
        self.logo_label.grid(row=0, column=0, padx=20, pady=(20, 10))

        # Introduction text
        self.intro_label = ctk.CTkLabel(
            self.sidebar_frame, 
            text=("This tool allows you to generate PowerPoint presentations "
                  "from Excel data using a predefined template. Upload your "
                  "Excel file and PowerPoint template, select the output folder, "
                  "and start the processing."), 
            font=ctk.CTkFont(size=16), 
            wraplength=270,  # Adjust the wrap length to fit your sidebar
            anchor="w",
            justify="center"
        )
        self.intro_label.grid(row=1, column=0, padx=20, pady=(10, 20), sticky="w")

        # Sidebar elements
        self.appearance_mode_label = ctk.CTkLabel(self.sidebar_frame, text="Appearance Mode:", anchor="w")
        self.appearance_mode_label.grid(row=2, column=0, padx=20, pady=(10, 0))
        
        self.appearance_mode_optionemenu = ctk.CTkOptionMenu(
            self.sidebar_frame, 
            values=["Light", "Dark"], 
            command=self.change_appearance_mode_event
        )
        self.appearance_mode_optionemenu.grid(row=3, column=0, padx=20, pady=(10, 10))
        
        self.scaling_label = ctk.CTkLabel(self.sidebar_frame, text="UI Scaling:", anchor="w")
        self.scaling_label.grid(row=4, column=0, padx=20, pady=(10, 0))
        
        self.scaling_optionemenu = ctk.CTkOptionMenu(
            self.sidebar_frame, 
            values=["80%", "90%", "100%", "110%", "120%"], 
            command=self.change_scaling_event
        )
        self.scaling_optionemenu.grid(row=5, column=0, padx=20, pady=(10, 20))

        # Create tabs
        self.tabview = ctk.CTkTabview(self, width=800)
        self.tabview.grid(row=0, column=1, columnspan=3, padx=(20, 20), pady=(20, 20), sticky="nsew")
        self.tabview.add("Input Data")
        self.tabview.add("Processing")
        self.tabview.add("Results")

        # Input Data Tab
        self.input_tab = self.tabview.tab("Input Data")
        self.input_tab.grid_columnconfigure(0, weight=1)
        self.input_tab.grid_rowconfigure(9, weight=1)

        # Input fields for Excel and PowerPoint template
        self.excel_file_var = tk.StringVar()
        self.pptx_template_file_var = tk.StringVar()
        self.output_folder_var = tk.StringVar()
        self.start_row_var = tk.StringVar()  # Removed default value
        self.end_row_var = tk.StringVar()

        self.create_input_fields()

        # Processing Tab
        self.processing_tab = self.tabview.tab("Processing")
        self.processing_tab.grid_columnconfigure(0, weight=1)
        self.processing_tab.grid_rowconfigure(1, weight=1)

        self.process_button = ctk.CTkButton(self.processing_tab, text="Generate PPT", command=self.run_processing, corner_radius=8)
        self.process_button.grid(row=1, column=0, padx=20, pady=20)

        # Results Tab
        self.results_tab = self.tabview.tab("Results")
        self.results_tab.grid_columnconfigure(0, weight=1)
        self.results_tab.grid_rowconfigure(0, weight=1)

        self.results_textbox = ctk.CTkTextbox(self.results_tab, width=300, wrap="word", corner_radius=10)
        self.results_textbox.grid(row=0, column=0, padx=20, pady=20, sticky="nsew")

        # Set default values
        self.appearance_mode_optionemenu.set("Dark")
        self.scaling_optionemenu.set("100%")

    def create_input_fields(self):
        ctk.CTkLabel(self.input_tab, text="Upload Excel File with Data", anchor="w").grid(row=0, column=0, padx=20, pady=10, sticky="w")
        ctk.CTkEntry(self.input_tab, textvariable=self.excel_file_var, placeholder_text="Select Excel file", width=500).grid(row=1, column=0, padx=20, pady=10, sticky="ew")
        ctk.CTkButton(self.input_tab, text="Browse Excel", command=self.browse_excel_file).grid(row=1, column=1, padx=20, pady=10)

        ctk.CTkLabel(self.input_tab, text="Upload PowerPoint Template", anchor="w").grid(row=2, column=0, padx=20, pady=10, sticky="w")
        ctk.CTkEntry(self.input_tab, textvariable=self.pptx_template_file_var, placeholder_text="Select PPTX file", width=500).grid(row=3, column=0, padx=20, pady=10, sticky="ew")
        ctk.CTkButton(self.input_tab, text="Browse PPTX", command=self.browse_pptx_template_file).grid(row=3, column=1, padx=20, pady=10)

        ctk.CTkLabel(self.input_tab, text="Output Folder Path", anchor="w").grid(row=4, column=0, padx=20, pady=10, sticky="w")
        ctk.CTkEntry(self.input_tab, textvariable=self.output_folder_var, placeholder_text="Select Output folder", width=500).grid(row=5, column=0, padx=20, pady=10, sticky="ew")
        ctk.CTkButton(self.input_tab, text="Browse Folder", command=self.browse_output_folder).grid(row=5, column=1, padx=20, pady=10)

        ctk.CTkLabel(self.input_tab, text="Starting Row (Excel)", anchor="w").grid(row=6, column=0, padx=20, pady=10, sticky="w")
        ctk.CTkEntry(self.input_tab, textvariable=self.start_row_var, width=500).grid(row=7, column=0, padx=20, pady=10, sticky="ew")

        ctk.CTkLabel(self.input_tab, text="Ending Row (Excel)", anchor="w").grid(row=8, column=0, padx=20, pady=10, sticky="w")
        ctk.CTkEntry(self.input_tab, textvariable=self.end_row_var, width=500).grid(row=9, column=0, padx=20, pady=10, sticky="ew")

    def browse_excel_file(self):
        filename = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx")])
        self.excel_file_var.set(filename)

    def browse_pptx_template_file(self):
        filename = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        self.pptx_template_file_var.set(filename)

    def browse_output_folder(self):
        foldername = filedialog.askdirectory()
        self.output_folder_var.set(foldername)

    def run_processing(self):
        try:
            # Collect inputs
            excel_file = self.excel_file_var.get()
            pptx_template_file = self.pptx_template_file_var.get()
            output_folder = self.output_folder_var.get()
            start_row = int(self.start_row_var.get())-1
            end_row = int(self.end_row_var.get()) if self.end_row_var.get() else None

            # Generate PPT from Excel data
            output_folder_path = self.generate_ppt_from_excel(excel_file, pptx_template_file, output_folder, start_row, end_row)

            # Show success message and update results tab
            messagebox.showinfo("Success", f"PPT files generated successfully in {output_folder_path}")
            self.results_textbox.insert("0.0", f"PPT files generated successfully in {output_folder_path}\n")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def generate_ppt_from_excel(self, excel_file, pptx_template, output_path, start_row=0, end_row=None):
        df = pd.read_excel(excel_file, sheet_name='Summaries')
        if end_row is None:
            end_row = len(df)

        for index, row in df.iloc[start_row:end_row].iterrows():
            prs = Presentation(pptx_template)

            shapes_data = {}
            for column in df.columns:
                if column in row:
                    value = row[column]
                    if pd.notna(value):
                        if column == 'Duckers Solution':
                            shapes_data[column] = self.replace_bullet_points(value)
                        else:
                            shapes_data[column] = value
                    else:
                        shapes_data[column] = ""

            for slide in prs.slides:
                self.update_shapes_with_excel_data(slide, shapes_data)

            output_folder = os.path.join(output_path, f"{row.get('Folder Name', 'Default')}")
            os.makedirs(output_folder, exist_ok=True)

            output_pptx_file = os.path.join(output_folder, f"{row.get('Case Study Name', 'Slide')}.pptx")

            prs.save(output_pptx_file)
            print(f"Saved {output_pptx_file} in {output_folder}")

        return output_path

    def replace_bullet_points(self, text):
        # Replace '*' with '•'
        return text.replace('*', '• ') if text is not None else ""

    def update_shapes_with_excel_data(self, slide, shapes_data):
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_name = shape.name
                if shape_name in shapes_data:
                    shape.text_frame.text = shapes_data[shape_name] if shapes_data[shape_name] is not None else ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            paragraph.alignment = PP_ALIGN.LEFT

    def change_appearance_mode_event(self, mode):
        ctk.set_appearance_mode(mode)

    def change_scaling_event(self, scale):
        scaling = int(scale.replace('%', '')) / 100
        ctk.set_widget_scaling(scaling)

if __name__ == "__main__":
    app = App()
    app.mainloop()

